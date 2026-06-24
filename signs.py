from __future__ import annotations

import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


EXPECTED_HEADER = [
	"#",
	"Assignment Code",
	"Judge Position",
	"Team",
	"Assigned Judge",
]


class ParseError(ValueError):
	def __init__(self, line_number: int, message: str) -> None:
		super().__init__(f"Line {line_number}: {message}")
		self.line_number = line_number
		self.message = message


@dataclass(frozen=True)
class JudgeRecord:
	competition_name: str | None
	session_name: str
	station_number: str
	assignment_code: str
	judge_position: str
	assigned_judge: str | None


def main(argv: list[str] | None = None) -> int:
	args = list(sys.argv[1:] if argv is None else argv)
	if len(args) != 1:
		print("Usage: python signs.py <input-file>", file=sys.stderr)
		return 1

	input_path = Path(args[0]).expanduser()
	if not input_path.is_file():
		print(f"Input file not found: {input_path}", file=sys.stderr)
		return 1

	try:
		records = parse_input_file(input_path)
		if not records:
			raise ValueError("The input file did not contain any judge assignments.")
		output_path = build_presentation(records, input_path)
	except (OSError, ParseError, ValueError) as exc:
		print(str(exc), file=sys.stderr)
		return 1

	print(f"Created {output_path}")
	return 0


def parse_input_file(input_path: Path) -> list[JudgeRecord]:
	text = input_path.read_text(encoding="utf-8-sig")
	lines = text.splitlines()

	competition_name: str | None = None
	session_name: str | None = None
	station_number: str | None = None
	expecting_header = False
	records: list[JudgeRecord] = []

	index = 0
	while index < len(lines):
		line_number = index + 1
		raw_line = lines[index]
		line = raw_line.strip()
		if not line:
			index += 1
			continue

		if line.startswith("Competition: "):
			if competition_name is not None:
				raise ParseError(line_number, "Competition was already defined.")
			if session_name is not None or station_number is not None or records:
				raise ParseError(line_number, "Competition must appear before session and station data.")
			competition_name = line[len("Competition: ") :].strip()
			if not competition_name:
				raise ParseError(line_number, "Competition name is empty.")
			index += 1
			continue

		if line.startswith("Judges: "):
			session_name = line[len("Judges: ") :].strip()
			if not session_name:
				raise ParseError(line_number, "Session name is empty.")
			station_number = None
			expecting_header = False
			index += 1
			continue

		if line.startswith("Station "):
			if session_name is None:
				raise ParseError(line_number, "Station found before a Judges session header.")
			station_number = line[len("Station ") :].strip()
			if not station_number:
				raise ParseError(line_number, "Station number is empty.")
			expecting_header = True
			index += 1
			continue

		if "\t" not in raw_line:
			if session_name is not None and station_number is not None and not expecting_header:
				index += 1
				continue
			raise ParseError(line_number, "Unexpected line content.")

		columns = raw_line.split("\t")
		if columns == EXPECTED_HEADER:
			if session_name is None:
				raise ParseError(line_number, "Header found before a Judges session header.")
			if station_number is None:
				raise ParseError(line_number, "Header found before a Station header.")
			expecting_header = False
			index += 1
			continue

		if expecting_header:
			raise ParseError(line_number, "Expected the tab-delimited judge header after the Station line.")

		if session_name is None:
			raise ParseError(line_number, "Judge row found before a Judges session header.")
		if station_number is None:
			raise ParseError(line_number, "Judge row found before a Station header.")

		# Support alternate exports where Assigned Judge is moved to the next line.
		if len(columns) >= 5 and not columns[4].strip() and index + 1 < len(lines):
			next_raw_line = lines[index + 1]
			next_line = next_raw_line.strip()
			if (
				next_line
				and "\t" not in next_raw_line
				and not next_line.startswith("Competition: ")
				and not next_line.startswith("Judges: ")
				and not next_line.startswith("Station ")
			):
				columns = columns.copy()
				columns[4] = next_line
				index += 1

		record = parse_judge_row(line_number, columns, competition_name, session_name, station_number)
		records.append(record)
		index += 1

	if expecting_header:
		raise ParseError(len(lines) or 1, "The file ended before the expected judge header row.")
	if session_name is None:
		raise ValueError("The input file did not contain any Judges sections.")

	return records


def parse_judge_row(
	line_number: int,
	columns: list[str],
	competition_name: str | None,
	session_name: str,
	station_number: str,
) -> JudgeRecord:
	if len(columns) == 4:
		columns = columns + [""]
	if len(columns) != 5:
		raise ParseError(
			line_number,
			"Judge row must contain 5 tab-delimited columns: #, Assignment Code, Judge Position, Team, Assigned Judge.",
		)

	assignment_code = columns[1].strip()
	judge_position = columns[2].strip()
	assigned_judge = normalize_judge_name(columns[4])

	if not assignment_code:
		raise ParseError(line_number, "Assignment Code is required.")
	if not judge_position:
		raise ParseError(line_number, "Judge Position is required.")

	return JudgeRecord(
		competition_name=competition_name,
		session_name=session_name,
		station_number=station_number,
		assignment_code=assignment_code,
		judge_position=judge_position,
		assigned_judge=assigned_judge,
	)


def normalize_judge_name(value: str) -> str | None:
	judge_name = value.strip()
	if not judge_name or judge_name == "-":
		return None
	return judge_name


def build_presentation(records: list[JudgeRecord], input_path: Path) -> Path:
	presentation = Presentation()

	while presentation.slides:
		slide_id = presentation.slides._sldIdLst[0]  # type: ignore[attr-defined]
		presentation.part.drop_rel(slide_id.rId)
		del presentation.slides._sldIdLst[0]  # type: ignore[attr-defined]

	slide_layout = presentation.slide_layouts[5]
	for record in records:
		add_judge_slide(presentation, slide_layout, record)

	output_path = input_path.with_suffix(".pptx")
	presentation.save(output_path)
	return output_path


def add_judge_slide(presentation: Presentation, slide_layout, record: JudgeRecord) -> None:
	slide = presentation.slides.add_slide(slide_layout)

	title = slide.shapes.title
	if title is not None:
		title.text = record.competition_name or ""

	body_width = title.width if title is not None else presentation.slide_width - Inches(1.0)
	body_height = Inches(4.0)
	body_left = title.left if title is not None else int((presentation.slide_width - body_width) / 2)
	body_top = int((presentation.slide_height - body_height) / 2)
	body_shape = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
	text_frame = body_shape.text_frame
	text_frame.clear()
	text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

	body_lines: list[str] = []
	if record.assigned_judge:
		body_lines.append(record.assigned_judge)
	body_lines.extend(
		[
			record.judge_position,
			record.session_name,
			f"Station {record.station_number}",
			record.assignment_code,
		]
	)

	for index, line in enumerate(body_lines):
		paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
		paragraph.text = line
		paragraph.alignment = PP_ALIGN.CENTER
		paragraph.font.size = Pt(32)


if __name__ == "__main__":
	raise SystemExit(main())
